import tkinter as tk
from tkinter import ttk, filedialog, messagebox
import os
import json
import re
import tempfile
import csv

# External libraries
import PyPDF2
import docx
import pandas as pd
from pptx import Presentation
import easyocr
from pdf2image import convert_from_path

# Initialize EasyOCR
ocr_reader = easyocr.Reader(['en'], gpu=False)  # Set gpu=True if available

OUTPUT_CSV = "output_results.csv"

# Regex patterns
ID_PATTERN = r'(([0-9]{2})(0|1)([0-9])([0-3])([0-9])([-.,#$%& ]?)([0-9]{4})([-.,#$%& ]?)([0-1][8]([-.,#$%& ]?)[0-9]))'
PASSPORT_PATTERN = r'\b[A|D|M|T][0-9]{8}\b'
CELLPHONE_PATTERN = r'(?:\+|00)?(?:0|27)?[5-9][0-9](?: )?[0-9]{3}(?: )?[0-9]{4}'
LANDLINE_PATTERN = r'(?:0|\+27)[0-9]{2}[-. ()]?[0-9]{3}[-. ()]?[0-9]{4}'
EMAIL_PATTERN = r'[\w\.-]+@(?:[\w-]+\.)+[\w-]{2,4}'
IPV4_IPV6_PATTERN = r'\b(?:(?:25[0-5]|2[0-4]\d|1\d{2}|[1-9]?\d)(?:\.(?!$)|$)){4}\b|\b(?:[A-Fa-f0-9]{1,4}:){1,7}[A-Fa-f0-9]{1,4}\b'
countries = [
    "Afghanistan", "Albania", "Algeria", "Andorra", "Angola", "Antigua and Barbuda", "Argentina", "Armenia",
    "Australia", "Austria", "Azerbaijan", "Bahamas", "Bahrain", "Bangladesh", "Barbados", "Belarus",
    "Belgium", "Belize", "Benin", "Bhutan", "Bolivia", "Bosnia and Herzegovina", "Botswana", "Brazil",
    "Brunei", "Bulgaria", "Burkina Faso", "Burundi", "Cabo Verde", "Cambodia", "Cameroon", "Canada",
    "Central African Republic", "Chad", "Chile", "China", "Colombia", "Comoros", "Congo (Congo-Brazzaville)",
    "Costa Rica", "Croatia", "Cuba", "Cyprus", "Czech Republic", "Democratic Republic of the Congo", "Denmark",
    "Djibouti", "Dominica", "Dominican Republic", "Ecuador", "Egypt", "El Salvador", "Equatorial Guinea",
    "Eritrea", "Estonia", "Eswatini", "Ethiopia", "Fiji", "Finland", "France", "Gabon", "Gambia", "Georgia",
    "Germany", "Ghana", "Greece", "Grenada", "Guatemala", "Guinea", "Guinea-Bissau", "Guyana", "Haiti",
    "Honduras", "Hungary", "Iceland", "India", "Indonesia", "Iran", "Iraq", "Ireland", "Israel", "Italy",
    "Ivory Coast", "Jamaica", "Japan", "Jordan", "Kazakhstan", "Kenya", "Kiribati", "Kuwait", "Kyrgyzstan",
    "Laos", "Latvia", "Lebanon", "Lesotho", "Liberia", "Libya", "Liechtenstein", "Lithuania", "Luxembourg",
    "Madagascar", "Malawi", "Malaysia", "Maldives", "Mali", "Malta", "Marshall Islands", "Mauritania",
    "Mauritius", "Mexico", "Micronesia", "Moldova", "Monaco", "Mongolia", "Montenegro", "Morocco",
    "Mozambique", "Myanmar", "Namibia", "Nauru", "Nepal", "Netherlands", "New Zealand", "Nicaragua", "Niger",
    "Nigeria", "North Korea", "North Macedonia", "Norway", "Oman", "Pakistan", "Palau", "Palestine",
    "Panama", "Papua New Guinea", "Paraguay", "Peru", "Philippines", "Poland", "Portugal", "Qatar",
    "Republic of the Congo", "Romania", "Russia", "Rwanda", "Saint Kitts and Nevis", "Saint Lucia",
    "Saint Vincent and the Grenadines", "Samoa", "San Marino", "Sao Tome and Principe", "Saudi Arabia",
    "Senegal", "Serbia", "Seychelles", "Sierra Leone", "Singapore", "Slovakia", "Slovenia", "Solomon Islands",
    "Somalia", "South Africa", "South Korea", "South Sudan", "Spain", "Sri Lanka", "Sudan", "Suriname",
    "Sweden", "Switzerland", "Syria", "Taiwan", "Tajikistan", "Tanzania", "Thailand", "Timor-Leste", "Togo",
    "Tonga", "Trinidad and Tobago", "Tunisia", "Turkey", "Turkmenistan", "Tuvalu", "Uganda", "Ukraine",
    "United Arab Emirates", "United Kingdom", "United States", "Uruguay", "Uzbekistan", "Vanuatu",
    "Vatican City", "Venezuela", "Vietnam", "Yemen", "Zambia", "Zimbabwe"
]
GENDER_PATTERN = r'(?:m|M|male|Male|f|F|female|Female|FEMALE|MALE|Not prefer to say)\b'

race_keywords = [
    "white", "caucasian", "black", "african american", "afro-american",
    "asian", "east asian", "south asian", "southeast asian",
    "hispanic", "latino", "latina",
    "native american", "indigenous", "first nations", "inuit",
    "middle eastern", "arab", "persian",
    "pacific islander", "polynesian", "maori"
]

def clear_table():
    for row in results_table.get_children():
        results_table.delete(row)

def insert_result_into_table(file_name, counts):
    if counts is None:
        results_table.insert("", tk.END, values=(file_name, "Skipped", "-", "-", "-", "-", "-", "-", "-", "-"))
    else:
        results_table.insert(
            "",
            tk.END,
            values=(
                file_name,
                counts["ID Number"],
                counts["Passport Number"],
                counts["Cellphone Number"],
                counts["Landline Number"],
                counts["Email"],
                counts["IP Address"],
                counts["Country"],
                counts["Gender"],
                counts["Race"],
            )
        )


def get_file_extension(file_path):
    return os.path.splitext(file_path)[1].lower()


def extract_text_with_easyocr_from_pdf(pdf_path):
    text = ''
    with tempfile.TemporaryDirectory() as path:
        images = convert_from_path(pdf_path, output_folder=path)
        for i, image in enumerate(images):
            ocr_result = ocr_reader.readtext(image)
            page_text = ' '.join([item[1] for item in ocr_result])
            text += f'\n\n--- OCR Page {i + 1} ---\n{page_text}'
    return text or ""


def extract_text_with_easyocr_from_image(image_path):
    ocr_result = ocr_reader.readtext(image_path)
    text = ' '.join([item[1] for item in ocr_result])
    return text or ""


def search_patterns_in_text(text):
    # Normalize text
    text_lower = text.lower()

    # SA ID numbers
    id_matches_raw = re.findall(ID_PATTERN, text)
    id_numbers = []
    for match in id_matches_raw:
        raw_id = match[0]
        clean_id = re.sub(r'[-.,#$%&\s]', '', raw_id)
        id_numbers.append(clean_id)

    # SA Passport numbers
    passport_matches = re.findall(PASSPORT_PATTERN, text)

    cellphone_numbers = re.findall(CELLPHONE_PATTERN, text)
    landline_numbers = re.findall(LANDLINE_PATTERN, text)
    email_matches = re.findall(EMAIL_PATTERN, text)
    ip_matches = re.findall(IPV4_IPV6_PATTERN, text)

    # Match countries
    matched_countries = []
    for country in countries:
        if country.lower() in text_lower:
            matched_countries.append(country)

    gender_matches = re.findall(GENDER_PATTERN, text)

    matched_races = []
    for race in race_keywords:
        if race.lower() in text_lower:
            matched_races.append(race)

    return id_numbers, passport_matches, cellphone_numbers, landline_numbers, email_matches, ip_matches, matched_countries,\
            gender_matches, matched_races


def read_file(file_path):
    ext = get_file_extension(file_path)
    text = ""

    try:
        if ext == '.txt':
            with open(file_path, 'r', encoding='utf-8') as f:
                text = f.read()

        elif ext == '.json':
            with open(file_path, 'r', encoding='utf-8') as f:
                data = json.load(f)
                text = json.dumps(data, indent=2)

        elif ext == '.pdf':
            with open(file_path, 'rb') as f:
                reader = PyPDF2.PdfReader(f)
                for page in reader.pages:
                    page_text = page.extract_text()
                    if page_text:
                        text += page_text

            if not text.strip():
                print(f"[INFO] No extractable text in PDF: {file_path}. Using OCR...")
                text = extract_text_with_easyocr_from_pdf(file_path)

        elif ext == '.docx':
            doc = docx.Document(file_path)
            text = '\n'.join([para.text for para in doc.paragraphs])

        elif ext in ['.xlsx', '.xls']:
            df = pd.read_excel(file_path)
            text = df.head().to_string()

        elif ext == '.pptx':
            prs = Presentation(file_path)
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + '\n'

        elif ext in ['.jpg', '.jpeg', '.png', '.tiff']:
            print(f"[INFO] Image file detected: {file_path}. Using OCR...")
            text = extract_text_with_easyocr_from_image(file_path)

    except Exception as e:
        print(f"[ERROR] Failed to read file: {file_path}\n{e}")

    return text


def process_single_file(file_path, csv_writer):
    file_name = os.path.basename(file_path)
    text = read_file(file_path)

    if not text.strip():
        return file_name, None  # Indicate skipped file (no text)

    id_numbers, passport_numbers, cellphone_numbers, landline_numbers, email_addresses, ip_addresses, matched_countries, \
        gender_matches, matched_races = search_patterns_in_text(text)

    ids_str = "; ".join(sorted(set(id_numbers))) if id_numbers else ""
    passports_str = "; ".join(sorted(set(passport_numbers))) if passport_numbers else ""
    cellphones_str = "; ".join(sorted(set(cellphone_numbers))) if cellphone_numbers else ""
    landlines_str = "; ".join(sorted(set(landline_numbers))) if landline_numbers else ""
    emails_str = "; ".join(sorted(set(email_addresses))) if email_addresses else ""
    ips_str = "; ".join(sorted(set(ip_addresses))) if ip_addresses else ""
    countries_str = "; ".join(sorted(set(matched_countries))) if matched_countries else ""
    genders_str = "; ".join(sorted(set(gender_matches))) if gender_matches else ""
    races_str = "; ".join(sorted(set(matched_races))) if matched_races else ""

    csv_writer.writerow(
        [file_name, ids_str, passports_str, cellphones_str, landlines_str, emails_str, ips_str, countries_str,
         genders_str, races_str])

    # Return counts for UI update
    counts = {
        "ID Number": len(id_numbers),
        "Passport Number": len(passport_numbers),
        "Cellphone Number": len(cellphone_numbers),
        "Landline Number": len(landline_numbers),
        "Email": len(email_addresses),
        "IP Address": len(ip_addresses),
        "Country": len(matched_countries),
        "Gender": len(gender_matches),
        "Race": len(matched_races)
    }
    return file_name, counts


def select_file():
    file_path = filedialog.askopenfilename()
    if file_path:
        with open(OUTPUT_CSV, mode='w', newline='', encoding='utf-8') as csvfile:
            csv_writer = csv.writer(csvfile)
            csv_writer.writerow(
                ["File Name", "ID Number", "Passport Number", "Cellphone Number", "Landline Number", "Email",
                 "IP Address", "Country", "Gender", "Race"])

            status_label.config(text=f"Processing: {os.path.basename(file_path)}")
            root.update_idletasks()

            file_name, counts = process_single_file(file_path, csv_writer)

            clear_table()
            insert_result_into_table(file_name, counts)

        status_label.config(text="Processing complete!")
        messagebox.showinfo("Done", f"Finished processing file.\nOutput saved to:\n{os.path.abspath(OUTPUT_CSV)}")


def select_folder():
    folder_path = filedialog.askdirectory()
    if folder_path:
        with open(OUTPUT_CSV, mode='w', newline='', encoding='utf-8') as csvfile:
            csv_writer = csv.writer(csvfile)
            csv_writer.writerow(
                ["File Name", "ID Number", "Passport Number", "Cellphone Number", "Landline Number", "Email",
                 "IP Address", "Country", "Gender", "Race"])

            status_label.config(text="")
            clear_table()

            for item in os.listdir(folder_path):
                full_path = os.path.join(folder_path, item)
                if os.path.isfile(full_path):
                    status_label.config(text=f"Processing: {item}")
                    root.update_idletasks()

                    file_name, counts = process_single_file(full_path, csv_writer)
                    insert_result_into_table(file_name, counts)

                    root.update_idletasks()

        status_label.config(text="Processing complete!")
        messagebox.showinfo("Done", f"Finished processing folder.\nOutput saved to:\n{os.path.abspath(OUTPUT_CSV)}")


root = tk.Tk()
root.title("File or Location Selector")
root.geometry("900x500")

header_label = tk.Label(root, text="IntelligENS PII Tools", font=("Lucida Handwriting", 20, "bold"), bg="#F3EA00", fg="black", pady=15)
header_label.pack(fill=tk.X)

file_button = tk.Button(root, text="Select File", command=select_file, width=20)
file_button.pack(pady=10)

folder_button = tk.Button(root, text="Select Folder", command=select_folder, width=20)
folder_button.pack(pady=10)

status_label = tk.Label(root, text="", fg="blue")
status_label.pack(pady=5)

# Treeview table
columns = ("File Name", "ID Number", "Passport Number", "Cellphone Number", "Landline Number",
           "Email", "IP Address", "Country", "Gender", "Race")

results_table = ttk.Treeview(root, columns=columns, show="headings", height=15)

for col in columns:
    results_table.heading(col, text=col)
    results_table.column(col, width=90, anchor=tk.CENTER)

results_table.pack(expand=True, fill=tk.BOTH, padx=10, pady=10)

root.mainloop()